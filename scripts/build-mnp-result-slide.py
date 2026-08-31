#!/usr/bin/env python3
"""Build a 16:9 Megafon-format result slide (MNP as-was / as-will + barriers)."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "slides" / "assets"
OUT = ROOT / "docs" / "slides" / "result-before-after.pptx"

GREEN = RGBColor(0x00, 0xB9, 0x56)
GREEN_BAR = RGBColor(0x28, 0xB9, 0x56)
PURPLE = RGBColor(0x73, 0x19, 0x82)
INK = RGBColor(0x28, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x8A, 0x8A, 0x8A)
RULE = RGBColor(0xEE, 0xEE, 0xEE)

FONT = "GT Walsheim Pro"


def set_run(run, text, *, size, bold=False, color=INK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Keep East-Asian / complex-script fallbacks aligned with Megafon theme.
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_textbox(slide, l, t, w, h, text, *, size, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    set_run(p.add_run(), text, size=size, bold=bold, color=color, font=font)
    return box


def add_card(slide, l, t, w, h, line_color, adj=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.background()
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.5)
    shape.line.dash_style = MSO_LINE_DASH_STYLE.SQUARE_DOT
    try:
        shape.adjustments[0] = adj
    except Exception:
        pass
    return shape


def contain(path: Path, box_l, box_t, box_w, box_h):
    im = Image.open(path)
    ar = im.width / im.height
    box_ar = box_w / box_h
    if ar > box_ar:
        w = box_w
        h = box_w / ar
    else:
        h = box_h
        w = box_h * ar
    l = box_l + (box_w - w) / 2
    t = box_t + (box_h - h) / 2
    return int(l), int(t), int(w), int(h)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 13.333 in
    prs.slide_height = Emu(6858000)  # 7.5 in
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    add_textbox(
        slide,
        Inches(0.42),
        Inches(0.28),
        Inches(12.48),
        Inches(0.78),
        "ЦЕЛЕВОЙ СЦЕНАРИЙ MNP СНИМАЕТ 48 БАРЬЕРОВ ИЗ 68, В ТОМ ЧИСЛЕ 26 КРИТИЧЕСКИХ",
        size=17,
        bold=True,
        color=INK,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # As-was card
    add_card(slide, Inches(0.42), Inches(1.16), Inches(8.55), Inches(2.42), PURPLE, adj=0.07)
    add_textbox(slide, Inches(0.58), Inches(1.22), Inches(5.2), Inches(0.32), "КАК СЕЙЧАС", size=11, bold=True)
    add_textbox(
        slide,
        Inches(5.7),
        Inches(1.22),
        Inches(3.05),
        Inches(0.32),
        "ВАРИАНТ А · ПРЕЗЕНТАЦИЯ 1",
        size=9,
        bold=True,
        color=PURPLE,
        align=PP_ALIGN.RIGHT,
    )
    was = ASSETS / "as-was.png"
    l, t, w, h = contain(was, Inches(0.58), Inches(1.56), Inches(8.23), Inches(1.88))
    slide.shapes.add_picture(str(was), l, t, w, h)

    # As-will card
    add_card(slide, Inches(0.42), Inches(3.70), Inches(8.55), Inches(3.18), GREEN, adj=0.055)
    add_textbox(slide, Inches(0.58), Inches(3.76), Inches(5.2), Inches(0.32), "КАК СТАНЕТ", size=11, bold=True)
    add_textbox(
        slide,
        Inches(5.7),
        Inches(3.76),
        Inches(3.05),
        Inches(0.32),
        "ВАРИАНТ B · ЦЕЛЕВОЙ СЦЕНАРИЙ",
        size=9,
        bold=True,
        color=GREEN,
        align=PP_ALIGN.RIGHT,
    )
    will = ASSETS / "as-will.png"
    l, t, w, h = contain(will, Inches(0.58), Inches(4.12), Inches(8.23), Inches(2.60))
    slide.shapes.add_picture(str(will), l, t, w, h)

    # Metrics card
    add_card(slide, Inches(9.14), Inches(1.16), Inches(3.76), Inches(5.72), GREEN, adj=0.05)
    add_textbox(
        slide,
        Inches(9.32),
        Inches(1.24),
        Inches(3.4),
        Inches(0.28),
        "БАРЬЕРЫ",
        size=10,
        bold=True,
        color=GREEN,
    )

    facts = [
        ("68", "всего барьеров", INK),
        ("38", "критических", INK),
        ("48", "будет решено", GREEN),
        ("26", "критических будет решено", GREEN),
        ("20", "не будет решено", PURPLE),
    ]
    top = 1.56
    for i, (num, label, color) in enumerate(facts):
        y = top + i * 0.78
        add_textbox(slide, Inches(9.32), Inches(y), Inches(0.95), Inches(0.62), num, size=28, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(
            slide,
            Inches(10.28),
            Inches(y),
            Inches(2.4),
            Inches(0.62),
            label.upper(),
            size=11,
            bold=True,
            color=INK,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if i < len(facts) - 1:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9.32),
                Inches(y + 0.70),
                Inches(3.4),
                Emu(12700),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RULE
            line.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.32), Inches(5.62), Inches(3.4), Inches(1.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_BAR
    bar.line.fill.background()
    try:
        bar.adjustments[0] = 0.08
    except Exception:
        pass
    tf = bar.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(
        p.add_run(),
        "Запущен A/B-тест до 30.08. 31.08–02.09 — анализ и подведение итогов теста.",
        size=11,
        bold=True,
        color=WHITE,
    )

    add_textbox(
        slide,
        Inches(0.42),
        Inches(6.98),
        Inches(4.2),
        Inches(0.28),
        "МегаФон | 31.08.2026",
        size=8,
        bold=False,
        color=MUTED,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(path)
